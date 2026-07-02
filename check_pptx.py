from pptx import Presentation

prs = Presentation("ChanyoungHeo_Portfolio.pptx")
print(f"현재 슬라이드 수: {len(prs.slides)}")
for i, slide in enumerate(prs.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                texts.append(t[:50])
    preview = " | ".join(texts[:3])
    print(f"  슬라이드 {i+1}: {preview}")
